import os
import io
import pickle
import threading
import mimetypes
from pathlib import Path
from mcp.server.fastmcp import FastMCP
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaInMemoryUpload, MediaFileUpload

# ================= CONFIG =================


BASE_DIR = Path(__file__).parent.parent
CLIENT_SECRET_FILE = BASE_DIR / "oauth-client.json"
TOKEN_FILE = BASE_DIR / "token.pickle"
SCOPES = [
    "https://www.googleapis.com/auth/drive",
    "https://www.googleapis.com/auth/gmail.modify",
    "https://mail.google.com/",
]

mcp = FastMCP("Advanced Drive MCP")

# ================= AUTH =================

_local = threading.local()


def get_service():
    if not hasattr(_local, "drive_service") or _local.drive_service is None:
        _local.drive_service = get_drive_service()
    return _local.drive_service


def get_drive_service(creds=None):
    if creds:
        return build("drive", "v3", credentials=creds)

    creds_file = os.environ.get("MCP_CREDS_FILE")
    if creds_file and Path(creds_file).exists():
        try:
            with open(creds_file, "rb") as f:
                user_creds = pickle.load(f)
            if user_creds and user_creds.valid:
                return build("drive", "v3", credentials=user_creds)
            if user_creds and user_creds.expired and user_creds.refresh_token:
                user_creds.refresh(Request())
                return build("drive", "v3", credentials=user_creds)
        except Exception:
            pass

    local_creds = None
    if TOKEN_FILE.exists():
        try:
            with open(TOKEN_FILE, "rb") as f:
                local_creds = pickle.load(f)
        except Exception:
            TOKEN_FILE.unlink(missing_ok=True)

    if local_creds and local_creds.expired and local_creds.refresh_token:
        try:
            local_creds.refresh(Request())
            with open(TOKEN_FILE, "wb") as f:
                pickle.dump(local_creds, f)
        except Exception:
            TOKEN_FILE.unlink(missing_ok=True)
            local_creds = None

    if not local_creds or not local_creds.valid:
        flow = InstalledAppFlow.from_client_secrets_file(
            str(CLIENT_SECRET_FILE), SCOPES
        )
        local_creds = flow.run_local_server(port=0)
        with open(TOKEN_FILE, "wb") as f:
            pickle.dump(local_creds, f)

    return build("drive", "v3", credentials=local_creds)


# ================= CORE HELPERS =================

def find_item(name: str, parent_id: str = None, mime_type: str = None) -> list:
    """Find file/folder with fuzzy matching — exact, swap underscore/space, keyword contains."""
    drive_service = get_service()

    def _search(q):
        full_q = q + " and trashed=false"
        if parent_id:
            full_q += f" and '{parent_id}' in parents"
        if mime_type:
            full_q += f" and mimeType='{mime_type}'"
        results = drive_service.files().list(
            q=full_q,
            fields="files(id, name, mimeType, parents, size, modifiedTime)"
        ).execute()
        return results.get("files", [])

    # 1. Exact match
    files = _search(f"name='{name}'")
    if files:
        return files

    # 2. Underscore ↔ space swap
    alt_name = name.replace('_', ' ') if '_' in name else name.replace(' ', '_')
    if alt_name != name:
        files = _search(f"name='{alt_name}'")
        if files:
            return files

    # 3. Keyword contains search — progressively shorter
    base_name = os.path.splitext(name)[0]
    alt_base = os.path.splitext(alt_name)[0]
    keywords = [w for w in alt_base.replace('_', ' ').split() if len(w) > 2]

    if keywords:
        for length in range(len(keywords), 0, -1):
            keyword = ' '.join(keywords[:length])
            files = _search(f"name contains '{keyword}'")
            if files:
                original_clean = name.lower().replace('_', ' ').replace('.', ' ')
                def score(f):
                    fname = f['name'].lower().replace('_', ' ').replace('.', ' ')
                    orig_words = set(original_clean.split())
                    file_words = set(fname.split())
                    return len(orig_words & file_words)
                files.sort(key=score, reverse=True)
                return files

    return []


def resolve_path(path: str, create_missing: bool = False) -> str | None:
    """Resolve a folder path like 'Projects/2024/Reports' to a folder ID."""
    drive_service = get_service()
    if not path or path.strip("/") == "":
        return None

    parts = [p for p in path.strip("/").split("/") if p]
    parent_id = None

    for part in parts:
        folders = find_item(part, parent_id, mime_type="application/vnd.google-apps.folder")
        if folders:
            parent_id = folders[0]["id"]
        elif create_missing:
            metadata = {"name": part, "mimeType": "application/vnd.google-apps.folder"}
            if parent_id:
                metadata["parents"] = [parent_id]
            folder = drive_service.files().create(body=metadata, fields="id").execute()
            parent_id = folder["id"]
        else:
            return None

    return parent_id


def split_path(path: str) -> tuple[str, str]:
    """Split 'Projects/2024/notes.txt' → ('Projects/2024', 'notes.txt')"""
    parts = path.strip("/").split("/")
    return "/".join(parts[:-1]), parts[-1]


def _read_content(file_id: str, mime_type: str, file_name: str, drive_service) -> str:
    """Read content from any supported file type. Returns text string."""

    # ── Google Workspace ──────────────────────────────────────────────────────
    if mime_type == "application/vnd.google-apps.document":
        content = drive_service.files().export(fileId=file_id, mimeType="text/plain").execute()
        return content.decode("utf-8", errors="ignore")

    if mime_type == "application/vnd.google-apps.spreadsheet":
        content = drive_service.files().export(fileId=file_id, mimeType="text/csv").execute()
        return content.decode("utf-8", errors="ignore")

    if mime_type == "application/vnd.google-apps.presentation":
        content = drive_service.files().export(fileId=file_id, mimeType="text/plain").execute()
        return content.decode("utf-8", errors="ignore")

    # ── Plain text formats ────────────────────────────────────────────────────
    if mime_type in [
        "text/plain", "text/csv", "text/html", "text/markdown",
        "application/json", "application/xml", "text/xml",
        "text/javascript", "application/x-python", "text/x-python",
    ]:
        content = drive_service.files().get_media(fileId=file_id).execute()
        return content.decode("utf-8", errors="ignore")

    # ── PDF ───────────────────────────────────────────────────────────────────
    if mime_type == "application/pdf":
        try:
            import pdfplumber
            content = drive_service.files().get_media(fileId=file_id).execute()
            text = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            result = "\n".join(text)
            return result[:3000] if result else "Error: Could not extract text from PDF."
        except ImportError:
            return "Error: pdfplumber not installed. Run: pip install pdfplumber"

    # ── DOCX ──────────────────────────────────────────────────────────────────
    if mime_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ]:
        try:
            from docx import Document
            content = drive_service.files().get_media(fileId=file_id).execute()
            doc = Document(io.BytesIO(content))
            text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            return text[:3000] if text else "Error: Could not extract text from DOCX."
        except ImportError:
            return "Error: python-docx not installed. Run: pip install python-docx"

    # ── XLSX ──────────────────────────────────────────────────────────────────
    if mime_type in [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ]:
        try:
            import openpyxl
            content = drive_service.files().get_media(fileId=file_id).execute()
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True)
            text = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text.append(f"Sheet: {sheet}")
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join([str(c) if c is not None else "" for c in row])
                    if row_text.strip():
                        text.append(row_text)
            result = "\n".join(text)
            return result[:3000] if result else "Error: Could not extract content from XLSX."
        except ImportError:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

    # ── PPTX ──────────────────────────────────────────────────────────────────
    if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        try:
            from pptx import Presentation
            content = drive_service.files().get_media(fileId=file_id).execute()
            prs = Presentation(io.BytesIO(content))
            text = []
            for i, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {i}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text.append(f"  {shape.text.strip()}")
            result = "\n".join(text)
            return result[:3000] if result else "Error: Could not extract text from PPTX."
        except ImportError:
            return "Error: python-pptx not installed. Run: pip install python-pptx"

    # ── ZIP / binary — not readable as text ──────────────────────────────────
    return f"Error: File type '{mime_type}' cannot be read as text. Try downloading instead."


# ================= TOOLS =================

@mcp.tool()
def create_folder(path: str) -> dict:
    """
    Create a folder or nested folder structure.
    Examples: 'Projects', 'Projects/2024/Reports'
    """
    try:
        folder_id = resolve_path(path, create_missing=True)
        return {"status": "success", "message": f"Folder '{path}' created successfully.", "folder_id": folder_id}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def list_files(path: str = "", page_size: str = "50", page_token: str = "") -> dict:
    """
    List all files and folders at a given path with pagination.
    path: folder path (empty for root).
    page_size: files per page (default 50).
    page_token: token for next page from previous result.
    """
    drive_service = get_service()
    try:
        parent_id = resolve_path(path) if path else None
        if path and parent_id is None:
            return {"status": "error", "message": f"Folder '{path}' not found."}

        query = "trashed=false"
        if parent_id:
            query += f" and '{parent_id}' in parents"

        count = int(page_size) if page_size else 50
        kwargs = dict(
            q=query,
            pageSize=count,
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime)"
        )
        if page_token:
            kwargs["pageToken"] = page_token

        results = drive_service.files().list(**kwargs).execute()
        files = results.get("files", [])
        next_token = results.get("nextPageToken", "")

        if not files:
            return {"files": [], "message": "No files found.", "next_page_token": ""}

        return {
            "files": [
                {
                    "name": f["name"],
                    "id": f["id"],
                    "type": "folder" if f["mimeType"] == "application/vnd.google-apps.folder" else "file",
                    "size": f.get("size", "—"),
                    "modified": f.get("modifiedTime", "—")
                }
                for f in files
            ],
            "next_page_token": next_token,
            "has_more": bool(next_token)
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def create_text_file(path: str, content: str, overwrite: str = "", append: str = "") -> dict:
    """
    Create a text-based file in Google Drive with content.
    Supported types: .txt .md .csv .html .json .xml .py .js .yaml .yml .toml .ini .sh .bat
    Use overwrite='true' to replace. Use append='true' to add to existing.
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(path)
        ext = Path(file_name).suffix.lower()

        mime_map = {
            ".txt":  "text/plain",
            ".md":   "text/markdown",
            ".csv":  "text/csv",
            ".html": "text/html",
            ".htm":  "text/html",
            ".json": "application/json",
            ".xml":  "application/xml",
            ".py":   "text/x-python",
            ".js":   "text/javascript",
            ".yaml": "text/yaml",
            ".yml":  "text/yaml",
            ".toml": "text/plain",
            ".ini":  "text/plain",
            ".sh":   "text/plain",
            ".bat":  "text/plain",
            ".sql":  "text/plain",
            ".r":    "text/plain",
        }
        mime_type = mime_map.get(ext, "text/plain")

        if not ext:
            file_name += ".txt"
            mime_type = "text/plain"

        parent_id = resolve_path(folder_path, create_missing=True) if folder_path else None
        existing = find_item(file_name, parent_id)

        do_overwrite = overwrite.lower() == "true"
        do_append = append.lower() == "true"

        if existing:
            file_id = existing[0]["id"]
            if not do_overwrite and not do_append:
                return {"status": "exists", "message": f"'{file_name}' already exists. Use overwrite='true' or append='true'."}
            if do_append:
                existing_content = drive_service.files().get_media(fileId=file_id).execute()
                content = existing_content.decode("utf-8", errors="ignore") + "\n" + content
            drive_service.files().update(fileId=file_id, body={"trashed": True}).execute()

        media = MediaInMemoryUpload(content.encode("utf-8"), mimetype=mime_type)
        metadata = {"name": file_name}
        if parent_id:
            metadata["parents"] = [parent_id]

        file = drive_service.files().create(body=metadata, media_body=media, fields="id").execute()
        return {"status": "success", "message": f"File '{path}' created successfully.", "file_id": file["id"]}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def read_file(path: str) -> str:
    """
    Read content of a file from Google Drive.
    Supports: txt, md, csv, html, json, xml, py, js, pdf, docx, xlsx, pptx,
    Google Docs, Google Sheets, Google Slides.
    path: file name or path like 'Projects/notes.txt'
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, parent_id)

        if not files:
            return f"File '{path}' not found in Drive."

        file_id = files[0]["id"]
        mime_type = files[0].get("mimeType", "")

        return _read_content(file_id, mime_type, file_name, drive_service)

    except Exception as e:
        return f"Error reading file: {str(e)}"


@mcp.tool()
def delete_file(path: str, confirmed: str = "") -> dict:
    """
    Move a file or folder to trash.
    Call first without confirmed to get confirmation.
    Call with confirmed='true' after user confirms.
    """
    drive_service = get_service()
    try:
        folder_path, name = split_path(path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(name, parent_id)

        if not files:
            return {"status": "error", "message": f"'{path}' not found."}

        if confirmed != "true":
            return {
                "status": "confirmation_required",
                "message": f"Found '{files[0]['name']}'. Move to trash?",
                "file_name": files[0]["name"],
                "file_id": files[0]["id"],
            }

        drive_service.files().update(fileId=files[0]["id"], body={"trashed": True}).execute()
        return {"status": "success", "message": f"'{files[0]['name']}' moved to trash successfully."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def restore_file(path: str, confirmed: str = "") -> dict:
    """
    Restore a file or folder from trash.
    Call first without confirmed to get confirmation.
    Call with confirmed='true' after user confirms.
    """
    drive_service = get_service()
    try:
        folder_path, name = split_path(path)
        query = f"name='{name}' and trashed=true"
        results = drive_service.files().list(q=query, fields="files(id, name)").execute()
        files = results.get("files", [])

        if not files:
            # Try fuzzy match in trash
            alt_name = name.replace('_', ' ') if '_' in name else name.replace(' ', '_')
            query = f"name='{alt_name}' and trashed=true"
            results = drive_service.files().list(q=query, fields="files(id, name)").execute()
            files = results.get("files", [])

        if not files:
            return {"status": "error", "message": f"'{path}' not found in trash."}

        if confirmed != "true":
            return {
                "status": "confirmation_required",
                "message": f"Found '{files[0]['name']}' in trash. Restore it?",
                "file_name": files[0]["name"],
                "file_id": files[0]["id"],
            }

        drive_service.files().update(fileId=files[0]["id"], body={"trashed": False}).execute()
        return {"status": "success", "message": f"'{files[0]['name']}' restored successfully."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def move_file(source_path: str, destination_folder: str) -> dict:
    """
    Move a file or folder to another folder.
    Example: move_file('Projects/old.txt', 'Archive/2024')
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(source_path)
        src_parent = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, src_parent)

        if not files:
            return {"status": "error", "message": f"'{source_path}' not found."}

        file_id = files[0]["id"]
        old_parents = ",".join(files[0].get("parents", []))
        dest_parent = resolve_path(destination_folder, create_missing=True)

        drive_service.files().update(
            fileId=file_id,
            addParents=dest_parent,
            removeParents=old_parents,
            fields="id, parents"
        ).execute()
        return {"status": "success", "message": f"'{file_name}' moved to '{destination_folder}' successfully."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def rename_file(path: str, new_name: str) -> dict:
    """
    Rename a file or folder.
    Example: rename_file('Projects/old.txt', 'new.txt')
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, parent_id)

        if not files:
            return {"status": "error", "message": f"'{path}' not found."}

        drive_service.files().update(fileId=files[0]["id"], body={"name": new_name}).execute()
        return {"status": "success", "message": f"'{file_name}' renamed to '{new_name}' successfully."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def copy_file(source_path: str, destination_folder: str, new_name: str = "") -> dict:
    """
    Copy a file to another folder, optionally with a new name.
    Example: copy_file('Projects/report.txt', 'Archive/2024', 'report_backup.txt')
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(source_path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, parent_id)

        if not files:
            return {"status": "error", "message": f"'{source_path}' not found."}

        dest_parent = resolve_path(destination_folder, create_missing=True)
        body = {"name": new_name or file_name, "parents": [dest_parent]}
        copied = drive_service.files().copy(fileId=files[0]["id"], body=body, fields="id, name").execute()

        return {
            "status": "success",
            "message": f"'{file_name}' copied to '{destination_folder}/{copied['name']}' successfully.",
            "file_id": copied["id"]
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def search_files(query: str, path: str = "") -> list:
    """
    Search for files by keyword across Drive or within a folder.
    NEVER include file extension in query. Use short keywords only.
    Example: search_files('milestones') not search_files('AI_ML Milestones.docx')
    """
    drive_service = get_service()
    try:
        parent_id = resolve_path(path) if path else None
        search_query = f"name contains '{query}' and trashed=false"
        if parent_id:
            search_query += f" and '{parent_id}' in parents"

        results = drive_service.files().list(
            q=search_query,
            fields="files(id, name, mimeType, modifiedTime)"
        ).execute()

        files = results.get("files", [])
        if not files:
            return [{"message": f"No files found matching '{query}'."}]

        return [
            {
                "name": f["name"],
                "id": f["id"],
                "type": "folder" if f["mimeType"] == "application/vnd.google-apps.folder" else "file",
                "modified": f.get("modifiedTime", "—")
            }
            for f in files
        ]
    except Exception as e:
        return [{"status": "error", "message": str(e)}]


@mcp.tool()
def get_file_info(path: str) -> dict:
    """Get detailed info about a file or folder."""
    try:
        folder_path, file_name = split_path(path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, parent_id)

        if not files:
            return {"status": "error", "message": f"'{path}' not found."}

        f = files[0]
        return {
            "name": f["name"],
            "id": f["id"],
            "type": "folder" if f["mimeType"] == "application/vnd.google-apps.folder" else "file",
            "mimeType": f["mimeType"],
            "size": f.get("size", "—"),
            "modified": f.get("modifiedTime", "—"),
            "parents": f.get("parents", [])
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def upload_file(file_path: str, folder_name: str = "", new_name: str = "") -> dict:
    """
    Upload any file from local machine to Google Drive.
    file_path: full local path of the file.
    folder_name: Drive folder to upload into (optional).
    new_name: rename file in Drive (optional).
    Supports ALL file types — images, videos, docs, code, zip, etc.
    """
    drive_service = get_service()
    try:
        if not os.path.exists(file_path):
            return {"status": "error", "message": f"File not found: {file_path}"}

        original_name = os.path.basename(file_path)
        upload_name = new_name if new_name else original_name
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            mime_type = "application/octet-stream"

        file_metadata = {"name": upload_name}
        if folder_name:
            folder_id = resolve_path(folder_name)
            if folder_id:
                file_metadata["parents"] = [folder_id]

        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        uploaded = drive_service.files().create(
            body=file_metadata, media_body=media, fields="id, name, mimeType"
        ).execute()

        return {
            "status": "success",
            "message": f"'{upload_name}' uploaded to Drive successfully.",
            "id": uploaded.get("id"),
            "name": uploaded.get("name"),
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def download_file(path: str, download_folder: str = "") -> dict:
    """
    Download a file from Google Drive to local machine.
    path: file name or path like 'Projects/notes.txt'.
    download_folder: local save folder (optional, defaults to ~/Downloads).
    Google Docs/Sheets/Slides are exported as docx/xlsx/pptx automatically.
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, parent_id)

        if not files:
            return {"status": "error", "message": f"File '{path}' not found in Drive."}

        file_id = files[0]["id"]
        mime_type = files[0].get("mimeType", "")
        actual_name = files[0]["name"]

        export_map = {
            "application/vnd.google-apps.document": (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"),
            "application/vnd.google-apps.spreadsheet": (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"),
            "application/vnd.google-apps.presentation": (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"),
        }

        if not download_folder:
            download_folder = str(Path.home() / "Downloads")
        os.makedirs(download_folder, exist_ok=True)

        if mime_type in export_map:
            export_mime, ext = export_map[mime_type]
            if not actual_name.endswith(ext):
                actual_name += ext
            content = drive_service.files().export(fileId=file_id, mimeType=export_mime).execute()
        else:
            content = drive_service.files().get_media(fileId=file_id).execute()

        # Always overwrite — Drive is source of truth
        save_path = os.path.join(download_folder, actual_name)
        with open(save_path, "wb") as f:
            f.write(content)

        return {
            "status": "success",
            "message": f"'{actual_name}' downloaded successfully.",
            "saved_to": save_path,
            "attachment_path": save_path,
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def list_recent_files(limit: str = "10", page_token: str = "") -> dict:
    """
    List recently modified files in Google Drive with pagination.
    limit: number of files to return (default 10).
    page_token: token for next page from previous result.
    """
    drive_service = get_service()
    try:
        count = int(limit) if limit else 10
        kwargs = dict(
            q="trashed=false",
            orderBy="modifiedTime desc",
            pageSize=count,
            fields="nextPageToken, files(id, name, mimeType, size, modifiedTime)"
        )
        if page_token:
            kwargs["pageToken"] = page_token

        results = drive_service.files().list(**kwargs).execute()
        files = results.get("files", [])
        next_token = results.get("nextPageToken", "")

        if not files:
            return {"files": [], "message": "No recent files found.", "next_page_token": ""}

        return {
            "files": [
                {
                    "name": f["name"],
                    "id": f["id"],
                    "type": "folder" if f["mimeType"] == "application/vnd.google-apps.folder" else "file",
                    "size": f.get("size", "—"),
                    "modified": f.get("modifiedTime", "—")
                }
                for f in files
            ],
            "next_page_token": next_token,
            "has_more": bool(next_token)
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}


@mcp.tool()
def share_file(path: str, email: str, role: str = "reader") -> dict:
    """
    Share a file or folder with someone via email.
    role: 'reader', 'commenter', or 'writer'.
    """
    drive_service = get_service()
    try:
        folder_path, file_name = split_path(path)
        parent_id = resolve_path(folder_path) if folder_path else None
        files = find_item(file_name, parent_id)

        if not files:
            return {"status": "error", "message": f"File '{path}' not found in Drive."}

        file_id = files[0]["id"]
        if role not in ["reader", "commenter", "writer"]:
            role = "reader"

        drive_service.permissions().create(
            fileId=file_id,
            body={"type": "user", "role": role, "emailAddress": email},
            sendNotificationEmail=True,
            fields="id"
        ).execute()

        return {"status": "success", "message": f"'{file_name}' shared with {email} as {role} successfully."}
    except Exception as e:
        return {"status": "error", "message": str(e)}


if __name__ == "__main__":
    mcp.run()